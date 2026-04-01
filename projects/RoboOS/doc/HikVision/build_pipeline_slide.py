"""
Single slide: Task Description → Claude Command Skills → Code Generation Pipeline
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x1B, 0x2A)
PANEL    = RGBColor(0x1A, 0x2B, 0x3C)
BLUE     = RGBColor(0x00, 0x8B, 0xD8)
CYAN     = RGBColor(0x00, 0xD4, 0xFF)
GREEN    = RGBColor(0x06, 0xD6, 0xA0)
ORANGE   = RGBColor(0xFF, 0x95, 0x00)
PURPLE   = RGBColor(0xA0, 0x6C, 0xFF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xAA, 0xBB, 0xCC)
DGRAY    = RGBColor(0x55, 0x66, 0x77)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG

def rect(l, t, w, h, color, line=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line: s.line.color.rgb = color
    else:    s.line.fill.background()
    return s

def tb(text, l, t, w, h, size=12, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def add_line(tf, text, size=11, bold=False, color=LGRAY,
             align=PP_ALIGN.LEFT, italic=False):
    p = tf.add_paragraph(); p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

# ── Top bar ───────────────────────────────────────────────────────────────────
rect(0, 0, Inches(13.33), Inches(0.08), CYAN)
rect(0, Inches(7.42),     Inches(13.33), Inches(0.08), CYAN)

# ── Title ─────────────────────────────────────────────────────────────────────
rect(0, Inches(0.08), Inches(13.33), Inches(0.82), PANEL)
tb("任务描述书  →  Claude Command Skills  →  代码自动生成",
   Inches(0.3), Inches(0.1), Inches(10), Inches(0.55),
   size=22, bold=True, color=WHITE)
tb("AI-Assisted Robot Skill Development Pipeline",
   Inches(0.3), Inches(0.62), Inches(9), Inches(0.3),
   size=12, italic=True, color=CYAN)

# ════════════════════════════════════════════════════════════════════════════
# Three main columns: INPUT  →  COMMAND SKILLS  →  OUTPUT
# ════════════════════════════════════════════════════════════════════════════
TOP  = Inches(1.05)
H    = Inches(5.9)
COL1_L = Inches(0.2)
COL2_L = Inches(4.0)
COL3_L = Inches(8.25)
COL_W  = Inches(3.6)

# ── Column backgrounds ────────────────────────────────────────────────────────
rect(COL1_L, TOP, COL_W, H, PANEL)
rect(COL2_L, TOP, COL_W, H, PANEL)
rect(COL3_L, TOP, COL_W, H, PANEL)

# ── Column header strips ──────────────────────────────────────────────────────
rect(COL1_L, TOP, COL_W, Inches(0.38), BLUE)
rect(COL2_L, TOP, COL_W, Inches(0.38), PURPLE)
rect(COL3_L, TOP, COL_W, Inches(0.38), GREEN)

tb("① INPUT  任务描述书",
   COL1_L+Inches(0.1), TOP+Inches(0.04), COL_W-Inches(0.15), Inches(0.32),
   size=13, bold=True, color=WHITE)
tb("② COMMAND SKILLS  技能命令",
   COL2_L+Inches(0.1), TOP+Inches(0.04), COL_W-Inches(0.15), Inches(0.32),
   size=13, bold=True, color=WHITE)
tb("③ OUTPUT  生成代码",
   COL3_L+Inches(0.1), TOP+Inches(0.04), COL_W-Inches(0.15), Inches(0.32),
   size=13, bold=True, color=WHITE)

# ── COL 1: Input doc content ───────────────────────────────────────────────────
# Doc icon area
rect(COL1_L+Inches(0.15), TOP+Inches(0.5), Inches(3.3), Inches(1.35),
     RGBColor(0x0A,0x20,0x38))
tb("📄  海康任务描述-01号任务.docx",
   COL1_L+Inches(0.25), TOP+Inches(0.55), Inches(3.1), Inches(0.35),
   size=11, bold=True, color=CYAN)

box1 = slide.shapes.add_textbox(
    COL1_L+Inches(0.25), TOP+Inches(0.92), Inches(3.1), Inches(0.85))
tf1 = box1.text_frame; tf1.word_wrap = True
tf1.paragraphs[0].alignment = PP_ALIGN.LEFT
r = tf1.paragraphs[0].add_run()
r.text = "九步操作流程描述"; r.font.size=Pt(10); r.font.color.rgb=LGRAY
for line in ["20项原子技能分类","力控 / 对准 / 同步精度要求","GR2双臂机器人能力说明"]:
    add_line(tf1, f"  ·  {line}", size=10, color=LGRAY)

# What Claude does: analyze
rect(COL1_L+Inches(0.15), TOP+Inches(2.0), Inches(3.3), Inches(0.28),
     RGBColor(0x00,0x50,0x90))
tb("Claude 分析文档", COL1_L+Inches(0.25), TOP+Inches(2.02),
   Inches(3.0), Inches(0.24), size=10, bold=True, color=CYAN)

box1b = slide.shapes.add_textbox(
    COL1_L+Inches(0.15), TOP+Inches(2.35), Inches(3.3), Inches(2.9))
tf1b = box1b.text_frame; tf1b.word_wrap=True
first = True
for item in [
    ("提炼原子技能库", ORANGE, True),
    ("  P1–P4  感知技能", LGRAY, False),
    ("  M1–M4  运动技能", LGRAY, False),
    ("  G1–G4  抓取技能", LGRAY, False),
    ("  O1–O4  操作技能", LGRAY, False),
    ("  C1       协作技能", LGRAY, False),
    ("  I1        交互技能", LGRAY, False),
    ("", LGRAY, False),
    ("识别设计约束", ORANGE, True),
    ("  力控精度 0.5–2N", LGRAY, False),
    ("  对准精度 ±0.3mm", LGRAY, False),
    ("  双臂同步 <100ms", LGRAY, False),
]:
    p = tf1b.paragraphs[0] if first else tf1b.add_paragraph()
    first = False
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = item[0]
    r.font.size = Pt(10); r.font.bold = item[2]; r.font.color.rgb = item[1]

# ── COL 2: Command Skills ──────────────────────────────────────────────────────
cmd_items = [
    ("/new-skill",   "new-skill.md",
     "新建原子技能 MCP Tool",
     ["输入: 技能ID + 机器人文件夹",
      "读取 skill.py 现有结构",
      "按规范生成 @mcp.tool() 函数",
      "写入力控/精度约束注释"],
     ORANGE),
    ("/plan-step",   "plan-step.md",
     "生成步骤执行计划",
     ["输入: 步骤编号 (1–9)",
      "查找规范技能序列",
      "对比已实现 vs 缺失技能",
      "输出伪代码调用链"],
     CYAN),
    ("/check-skills","check-skills.md",
     "审计技能覆盖率",
     ["扫描 skill.py 已实现技能",
      "输出 20项覆盖矩阵 ✅/❌",
      "9步工作流就绪状态",
      "优先补全建议"],
     GREEN),
]

for ci, (cmd, fname, desc, bullets, color) in enumerate(cmd_items):
    ty = TOP + Inches(0.5 + ci * 1.75)
    rect(COL2_L+Inches(0.15), ty, Inches(3.3), Inches(1.6),
         RGBColor(0x15, 0x10, 0x30))
    rect(COL2_L+Inches(0.15), ty, Inches(3.3), Inches(0.3), color)

    tb(cmd, COL2_L+Inches(0.22), ty+Inches(0.02),
       Inches(1.6), Inches(0.26), size=12, bold=True, color=BG)
    tb(fname, COL2_L+Inches(1.9), ty+Inches(0.04),
       Inches(1.5), Inches(0.22), size=9, italic=True, color=BG, align=PP_ALIGN.RIGHT)
    tb(desc, COL2_L+Inches(0.22), ty+Inches(0.35),
       Inches(3.0), Inches(0.28), size=10, bold=True, color=color)

    bx = slide.shapes.add_textbox(
        COL2_L+Inches(0.22), ty+Inches(0.65), Inches(3.0), Inches(0.85))
    btf = bx.text_frame; btf.word_wrap=True
    btf.paragraphs[0].alignment = PP_ALIGN.LEFT
    r = btf.paragraphs[0].add_run()
    r.text = f"▸  {bullets[0]}"; r.font.size=Pt(9); r.font.color.rgb=LGRAY
    for b in bullets[1:]:
        add_line(btf, f"▸  {b}", size=9, color=LGRAY)

# How to invoke label
rect(COL2_L+Inches(0.15), TOP+Inches(5.55), Inches(3.3), Inches(0.22),
     RGBColor(0x20,0x10,0x40))
tb("Claude Code 中使用  /new-skill  /plan-step  /check-skills  调用",
   COL2_L+Inches(0.2), TOP+Inches(5.56), Inches(3.2), Inches(0.2),
   size=8, italic=True, color=PURPLE, align=PP_ALIGN.CENTER)

# ── COL 3: Generated Code ──────────────────────────────────────────────────────
output_items = [
    ("gr2_robot_local/skill.py",    "20项 @mcp.tool() 原子技能实现\n(P/M/G/O/C/I 全类别，含精度约束)", GREEN),
    ("gr2_robot_remote/skill.py",   "HTTP 远程模式同等技能\n(stateless_http, port 8000)", GREEN),
    ("master/scene/profile.yaml",   "气密检测工作台场景\ntransferBoxIn/Out · inspectionStation", BLUE),
    ("master/agents/prompts.py",    "注入 GR2 上下文 + 9步参考流程\n约束条件嵌入规划提示词", ORANGE),
    ("slaver/config.yaml",          "robot: gr2_robot · local · max_tools=5", LGRAY),
    ("slaver/tools/memory.py",      "新增 none 动作类型\n感知/等待类技能不触发场景更新", LGRAY),
]

for ri, (fname, desc, color) in enumerate(output_items):
    ty = TOP + Inches(0.5 + ri * 0.88)
    rect(COL3_L+Inches(0.15), ty, Inches(0.08), Inches(0.75), color)
    rect(COL3_L+Inches(0.28), ty, Inches(3.17), Inches(0.75),
         RGBColor(0x0A,0x20,0x18) if color==GREEN else RGBColor(0x0A,0x18,0x28)
         if color==BLUE else RGBColor(0x1C,0x12,0x04) if color==ORANGE
         else RGBColor(0x14,0x18,0x1C))
    tb(fname, COL3_L+Inches(0.35), ty+Inches(0.05),
       Inches(3.0), Inches(0.28), size=10, bold=True, color=color)
    tb(desc, COL3_L+Inches(0.35), ty+Inches(0.35),
       Inches(3.0), Inches(0.38), size=9, color=LGRAY, wrap=True)

# ── ARROWS between columns ────────────────────────────────────────────────────
for ax, label in [(Inches(3.82), "Claude\nAnalyzes"), (Inches(7.87), "Claude\nGenerates")]:
    rect(ax, Inches(3.6), Inches(0.16), Inches(0.35), RGBColor(0x33,0x44,0x55))
    rect(ax, Inches(3.6), Inches(0.16), Inches(0.35), BG)  # invisible spacer trick
    # Draw arrow as text
    tb("➜", ax - Inches(0.05), Inches(3.45), Inches(0.5), Inches(0.5),
       size=26, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    tb(label, ax - Inches(0.05), Inches(3.95), Inches(0.5), Inches(0.55),
       size=8, color=LGRAY, align=PP_ALIGN.CENTER)

# ── Bottom summary bar ────────────────────────────────────────────────────────
rect(0, Inches(7.08), Inches(13.33), Inches(0.34), RGBColor(0x10,0x22,0x35))
tb("📄 任务文档  →  🧠 Claude 理解  →  📝 Command Skills (.md)  →  ⚡ 一键生成  →  🤖 完整机器人技能代码",
   Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.28),
   size=11, color=CYAN, align=PP_ALIGN.CENTER, bold=True)

OUT = "/home/haoanw/workspace/RoboOS/doc/RoboOS_Pipeline_Slide.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
