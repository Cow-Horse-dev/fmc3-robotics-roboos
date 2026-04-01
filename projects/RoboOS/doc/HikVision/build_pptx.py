"""
Build a PowerPoint presentation for the Hikvision camera airtightness inspection
workflow adapted to Fourier GR2 + RoboOS.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
BG_CARD    = RGBColor(0x16, 0x21, 0x3E)   # card navy
ACCENT     = RGBColor(0x0F, 0x3C, 0x78)   # medium blue
HIGHLIGHT  = RGBColor(0x00, 0xB4, 0xD8)   # cyan
GREEN      = RGBColor(0x06, 0xD6, 0xA0)   # teal green
ORANGE     = RGBColor(0xFF, 0xA5, 0x00)   # amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW     = RGBColor(0xFF, 0xD1, 0x66)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helper functions ─────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    fill_bg(slide, BG_DARK)
    return slide


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None, accent_color=HIGHLIGHT):
    """Top banner with title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), ACCENT)
    add_rect(slide, 0, 0, Inches(0.08), Inches(1.1), accent_color)
    add_textbox(slide, title,
                Inches(0.25), Inches(0.12), Inches(12), Inches(0.6),
                font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.25), Inches(0.7), Inches(12), Inches(0.35),
                    font_size=14, color=HIGHLIGHT, italic=True)


def skill_tag(slide, text, l, t, cat_color=GREEN):
    w, h = Inches(1.55), Inches(0.3)
    add_rect(slide, l, t, w, h, cat_color)
    add_textbox(slide, text, l + Inches(0.05), t, w - Inches(0.05), h,
                font_size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_bg(slide, BG_DARK)

# Decorative top bar
add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), HIGHLIGHT)
add_rect(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), HIGHLIGHT)

# Main title box
add_rect(slide, Inches(1.0), Inches(1.5), Inches(11.33), Inches(2.5), ACCENT)
add_rect(slide, Inches(1.0), Inches(1.5), Inches(0.12), Inches(2.5), HIGHLIGHT)

add_textbox(slide, "海康摄像头气密检测",
            Inches(1.3), Inches(1.65), Inches(11), Inches(1.0),
            font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(slide, "RoboOS × Fourier GR2 双臂人形机器人  原子技能适配方案",
            Inches(1.3), Inches(2.65), Inches(11), Inches(0.6),
            font_size=20, color=HIGHLIGHT, align=PP_ALIGN.LEFT)

# Metadata row
for i, (label, val) in enumerate([
    ("机器人", "Fourier GR2（双臂 · 五指灵巧手）"),
    ("原子技能数", "20 项 / 7 类"),
    ("任务步骤", "9 步全流程"),
]):
    lx = Inches(1.3 + i * 3.7)
    add_rect(slide, lx, Inches(4.4), Inches(3.4), Inches(0.8), BG_CARD)
    add_textbox(slide, label, lx + Inches(0.1), Inches(4.43), Inches(3.2), Inches(0.3),
                font_size=11, color=HIGHLIGHT, bold=True)
    add_textbox(slide, val,   lx + Inches(0.1), Inches(4.72), Inches(3.2), Inches(0.32),
                font_size=13, color=WHITE)

add_textbox(slide, "RoboOS Adaptation  ·  Hikvision Camera Airtightness Inspection",
            Inches(1.3), Inches(6.8), Inches(10), Inches(0.4),
            font_size=11, color=LIGHT_GRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Architecture
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "RoboOS 系统架构", "Brain-Cerebellum 两级控制架构")

# Three columns: Master / Redis / Slaver
col_data = [
    (ACCENT,   "🧠  Master（大脑）",
     ["全局任务规划", "LLM 任务分解", "9步有序子任务", "Redis 任务分发", "机器人状态监控"]),
    (RGBColor(0x1B, 0x4F, 0x72), "⚡  Redis 通信层",
     ["roboos_to_gr2_robot", "{robot}_to_RoboOS", "Pub/Sub 消息队列", "场景状态共享", "心跳 & 结果上报"]),
    (RGBColor(0x0A, 0x3D, 0x2E), "🤖  Slaver（小脑）",
     ["接收子任务", "ToolMatcher 语义匹配", "ToolCallingAgent ReAct", "MCP 工具调用", "SceneMemory 场景更新"]),
]
for ci, (col_color, title, items) in enumerate(col_data):
    lx = Inches(0.4 + ci * 4.3)
    add_rect(slide, lx, Inches(1.25), Inches(4.1), Inches(5.8), col_color)
    add_rect(slide, lx, Inches(1.25), Inches(4.1), Inches(0.5), HIGHLIGHT if ci == 0 else (ORANGE if ci == 1 else GREEN))
    add_textbox(slide, title, lx + Inches(0.1), Inches(1.27), Inches(3.9), Inches(0.45),
                font_size=14, bold=True, color=BG_DARK)
    for ri, item in enumerate(items):
        add_textbox(slide, f"▸  {item}",
                    lx + Inches(0.15), Inches(1.9 + ri * 0.75), Inches(3.8), Inches(0.6),
                    font_size=13, color=WHITE)

# Arrows between columns
for ax in [Inches(4.52), Inches(8.82)]:
    add_textbox(slide, "⟷", ax, Inches(3.8), Inches(0.5), Inches(0.5),
                font_size=22, color=HIGHLIGHT, align=PP_ALIGN.CENTER, bold=True)

# Bottom note
add_rect(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3), ACCENT)
add_textbox(slide, "底层 LLM：RoboBrain2.0-7B  via vLLM  ·  MCP 工具协议  ·  FastMCP stdio / HTTP 双模式",
            Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
            font_size=11, color=HIGHLIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 9-Step Workflow Overview
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "九步作业流程总览", "海康摄像头气密检测 — 完整操作序列")

steps = [
    ("1", "取摄像头",      "从周转盒(In)取出摄像头",         GREEN),
    ("2", "摘镜头帽+检查", "右手摘帽 → 目视检查镜头",        HIGHLIGHT),
    ("3", "扫描二维码",    "转向扫码器 → 完成扫码",          ORANGE),
    ("4", "放入工装",      "镜头朝上放入检测台槽位",          GREEN),
    ("5", "启动检测",      "双指同时按绿色按钮 → 等待",       HIGHLIGHT),
    ("6", "取出+复检",     "右手取出摄像头 → 再次目检",       ORANGE),
    ("7", "盖镜头帽",      "左手取帽 → 盖回镜头",            GREEN),
    ("8", "码放入盒",      "镜头朝下 → 放入输出周转盒",       HIGHLIGHT),
    ("9", "查看结果",      "转身读取电脑屏幕气密结果",         ORANGE),
]

cols = 3
rows = 3
for idx, (num, name, desc, color) in enumerate(steps):
    col = idx % cols
    row = idx // cols
    lx = Inches(0.35 + col * 4.3)
    ty = Inches(1.25 + row * 1.9)
    add_rect(slide, lx, ty, Inches(4.1), Inches(1.72), BG_CARD)
    add_rect(slide, lx, ty, Inches(0.45), Inches(1.72), color)
    add_textbox(slide, num,  lx + Inches(0.05), ty + Inches(0.05), Inches(0.35), Inches(0.5),
                font_size=22, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, name, lx + Inches(0.55), ty + Inches(0.1), Inches(3.45), Inches(0.45),
                font_size=14, bold=True, color=color)
    add_textbox(slide, desc, lx + Inches(0.55), ty + Inches(0.55), Inches(3.45), Inches(1.0),
                font_size=12, color=LIGHT_GRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Atomic Skill Library
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "原子技能库", "20 项原子技能 × 7 类  —  可复用组件")

cats = [
    ("P  感知", HIGHLIGHT, [
        ("P1", "visual_localize",       "视觉定位目标"),
        ("P2", "visual_inspect",        "镜头外观检测"),
        ("P3", "qr_code_recognize",     "二维码识别"),
        ("P4", "read_screen_result",    "读取检测结果"),
    ]),
    ("M  运动", GREEN, [
        ("M1", "move_to_position",      "单臂/双臂移位"),
        ("M2", "bimanual_sync_move",    "双臂同步运动"),
        ("M3", "set_orientation",       "姿态旋转控制"),
        ("M4", "plan_path",             "无碰撞路径规划"),
    ]),
    ("G  抓取", ORANGE, [
        ("G1", "open_hand",             "五指张开"),
        ("G2", "precision_pinch",       "双指精确抓取"),
        ("G3", "force_controlled_grasp","力控抓取 0.5-2N"),
        ("G4", "lift_object",           "抬起物体"),
    ]),
    ("O  操作", YELLOW, [
        ("O1", "place_object",          "放置物体"),
        ("O2", "press_dual_buttons",    "双按钮同步按压"),
        ("O3", "lens_cap_operation",    "镜头帽插拔"),
        ("O4", "fine_align",            "精对准 ±0.3mm"),
    ]),
    ("C  协作", RGBColor(0xC7, 0x77, 0xFF), [
        ("C1", "hand_transfer",         "双手传递物体"),
    ]),
    ("I  交互", RGBColor(0xFF, 0x77, 0xA8), [
        ("I1", "wait_for_signal",       "等待外部信号"),
    ]),
]

col_w = Inches(2.1)
for ci, (cat_name, cat_color, skills) in enumerate(cats):
    col = ci % 4
    row = ci // 4
    lx = Inches(0.3 + col * 3.25)
    ty = Inches(1.25 + row * 2.9)
    box_h = Inches(0.4 + len(skills) * 0.58)
    add_rect(slide, lx, ty, Inches(3.1), box_h, BG_CARD)
    add_rect(slide, lx, ty, Inches(3.1), Inches(0.38), cat_color)
    add_textbox(slide, cat_name, lx + Inches(0.08), ty + Inches(0.04),
                Inches(3.0), Inches(0.3), font_size=12, bold=True, color=BG_DARK)
    for si, (sid, fname, desc) in enumerate(skills):
        sy = ty + Inches(0.45 + si * 0.57)
        add_rect(slide, lx + Inches(0.1), sy, Inches(0.35), Inches(0.32), cat_color)
        add_textbox(slide, sid,   lx + Inches(0.1), sy, Inches(0.35), Inches(0.32),
                    font_size=9, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, fname, lx + Inches(0.52), sy, Inches(1.55), Inches(0.32),
                    font_size=10, bold=True, color=cat_color)
        add_textbox(slide, desc,  lx + Inches(0.52), sy + Inches(0.16), Inches(2.5), Inches(0.32),
                    font_size=9,  color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES 5-13 — One slide per workflow step
# ════════════════════════════════════════════════════════════════════════════
step_details = [
    {
        "num": 1, "name": "取摄像头产品",
        "desc": "工人（机器人）左手从周转盒内双指拿取摄像头产品，镜头帽朝上。",
        "seq":  "P1 → M4 → M1 → G1 → G2 → G3 → G4",
        "skills": [
            ("P1", "visual_localize", "定位摄像头在周转盒中的位置", HIGHLIGHT),
            ("M4", "plan_path",       "规划无碰撞路径至抓取点",    GREEN),
            ("M1", "move_to_position","左臂移动到抓取点上方",      GREEN),
            ("G1", "open_hand",       "张开左手五指",               ORANGE),
            ("G2", "precision_pinch", "双指精确定位摄像头",         ORANGE),
            ("G3", "force_controlled_grasp", "力控抓取 1N（避免损伤）", ORANGE),
            ("G4", "lift_object",     "抬起摄像头离开周转盒",       ORANGE),
        ],
        "notes": ["抓取力 0.5–2N，避免损伤镜头", "P1 视觉伺服补偿位置误差"],
        "color": GREEN,
    },
    {
        "num": 2, "name": "摘掉镜头帽，检查产品外观",
        "desc": "右手从摄像头上摘掉镜头帽，放置在气密检测台指定区域；目视检查镜头是否受损。",
        "seq":  "C1 → P1 → M1 → G2 → O3 → O1 → M3 → P2",
        "skills": [
            ("C1", "hand_transfer",      "左→右手传递摄像头",         RGBColor(0xC7,0x77,0xFF)),
            ("P1", "visual_localize",    "定位镜头帽位置",             HIGHLIGHT),
            ("M1", "move_to_position",   "左手移向镜头帽",             GREEN),
            ("G2", "precision_pinch",    "双指捏住镜头帽",             ORANGE),
            ("O3", "lens_cap_operation", "pull：拔出镜头帽",           YELLOW),
            ("O1", "place_object",       "将镜头帽放至 lens_cap_area", YELLOW),
            ("M3", "set_orientation",    "转摄像头镜头朝向人眼",        GREEN),
            ("P2", "visual_inspect",     "目视检测镜头有无划伤",        HIGHLIGHT),
        ],
        "notes": ["镜头帽放置区域固定，避免第七步位置交叉", "P2 检测结果写入 SceneMemory"],
        "color": HIGHLIGHT,
    },
    {
        "num": 3, "name": "扫描产品二维码",
        "desc": "双手转动摄像头，将机身侧面二维码转向扫码器方向，移动到扫码位完成扫码。",
        "seq":  "M2 → M3 → M1 → O4 → P3 → I1",
        "skills": [
            ("M2", "bimanual_sync_move",  "双手同步旋转摄像头",      GREEN),
            ("M3", "set_orientation",     "二维码面向扫码器方向",    GREEN),
            ("M1", "move_to_position",    "移动至扫码器下方",        GREEN),
            ("O4", "fine_align",          "精对准扫码器光束 ±0.3mm",YELLOW),
            ("P3", "qr_code_recognize",   "触发扫码器读取二维码",    HIGHLIGHT),
            ("I1", "wait_for_signal",     "等待扫码成功信号",        RGBColor(0xFF,0x77,0xA8)),
        ],
        "notes": ["双臂同步时间差 < 100ms", "I1 等待扫码成功确认后继续"],
        "color": ORANGE,
    },
    {
        "num": 4, "name": "水平放置于气密检测工装",
        "desc": "转动摄像头使镜头朝正上方，水平放置到工装中间的正方形槽位内。",
        "seq":  "M3 → P1 → M4 → M1 → O4 → O1",
        "skills": [
            ("M3", "set_orientation",  "旋转摄像头至镜头朝上",     GREEN),
            ("P1", "visual_localize",  "定位工装槽位精确坐标",     HIGHLIGHT),
            ("M4", "plan_path",        "规划路径避开工装边缘",     GREEN),
            ("M1", "move_to_position", "移动摄像头至槽位正上方",   GREEN),
            ("O4", "fine_align",       "精对准槽位 ±0.3mm",       YELLOW),
            ("O1", "place_object",     "水平放入 fixture_slot",   YELLOW),
        ],
        "notes": ["放置精度 ±0.3mm，防止工装卡死", "镜头必须朝正上方（M3 set_orientation='lens_up'）"],
        "color": GREEN,
    },
    {
        "num": 5, "name": "按下启动按钮，等待气密检测",
        "desc": "双手食指同时按下检测台最左侧和中间的两个绿色按钮，观察面板等待检测完成。",
        "seq":  "P1 → M2 → O2 → I1",
        "skills": [
            ("P1", "visual_localize",   "定位两个绿色按钮位置",       HIGHLIGHT),
            ("M2", "bimanual_sync_move", "双臂同步移向两个按钮",      GREEN),
            ("O2", "press_dual_buttons", "双食指同时按压（<100ms）", YELLOW),
            ("I1", "wait_for_signal",   "等待 airtightness_test_complete", RGBColor(0xFF,0x77,0xA8)),
        ],
        "notes": ["必须双指同时按下（时间差 < 100ms）", "检测时间约 30–60s，I1 timeout_s=120"],
        "color": HIGHLIGHT,
    },
    {
        "num": 6, "name": "从工装内取出产品，检查外观",
        "desc": "气密检测完成后，右手从工装槽位内取出摄像头，转镜头朝人眼再次目视检查。",
        "seq":  "M1 → G2 → G4 → M3 → P2",
        "skills": [
            ("M1", "move_to_position", "右手伸入工装槽位",      GREEN),
            ("G2", "precision_pinch",  "双指抓取摄像头",        ORANGE),
            ("G4", "lift_object",      "抬起取出摄像头",        ORANGE),
            ("M3", "set_orientation",  "转镜头朝向人眼方向",    GREEN),
            ("P2", "visual_inspect",   "目检镜头有无损伤",      HIGHLIGHT),
        ],
        "notes": ["与步骤2复检动作相同，确认检测中无损伤", "P2 结果记录至 SceneMemory"],
        "color": ORANGE,
    },
    {
        "num": 7, "name": "盖上镜头帽",
        "desc": "左手从指定区域拿取镜头帽，将镜头帽盖到右手中的摄像头镜头上。",
        "seq":  "P1 → M1 → G2 → O3 → O1",
        "skills": [
            ("P1", "visual_localize",    "定位 lens_cap_area 中的镜头帽", HIGHLIGHT),
            ("M1", "move_to_position",   "左手移动至镜头帽位置",          GREEN),
            ("G2", "precision_pinch",    "双指拿取镜头帽",                ORANGE),
            ("O3", "lens_cap_operation", "insert：对准镜头插入",          YELLOW),
            ("O1", "place_object",       "确认镜头帽盖合",                YELLOW),
        ],
        "notes": ["O3 insert 精度 ±0.3mm", "与步骤2中放帽位置对应"],
        "color": GREEN,
    },
    {
        "num": 8, "name": "整齐放入输出周转盒",
        "desc": "右手将摄像头放入右侧输出周转盒，镜头朝下，按从右到左、从上到下顺序排放。",
        "seq":  "P1 → M4 → M3 → O4 → O1",
        "skills": [
            ("P1", "visual_localize",  "定位 transferBoxOut 当前空槽位", HIGHLIGHT),
            ("M4", "plan_path",        "规划路径至目标槽位",              GREEN),
            ("M3", "set_orientation",  "旋转摄像头至镜头朝下",            GREEN),
            ("O4", "fine_align",       "对准盒内槽位 ±0.3mm",            YELLOW),
            ("O1", "place_object",     "放入 transferBoxOut",            YELLOW),
        ],
        "notes": ["排列顺序：右→左，上→下", "P1 需感知当前已放位置以确定下一空槽"],
        "color": HIGHLIGHT,
    },
    {
        "num": 9, "name": "查看屏幕气密检测结果",
        "desc": "转身查看左侧电脑显示屏上的气密检测结果，确认是否成功。",
        "seq":  "P1 → P4 → P2",
        "skills": [
            ("P1", "visual_localize",   "定位屏幕显示区域",          HIGHLIGHT),
            ("P4", "read_screen_result","OCR 读取检测结果文字",       HIGHLIGHT),
            ("P2", "visual_inspect",    "判断结果 PASS / FAIL",      HIGHLIGHT),
        ],
        "notes": ["结果写入 SceneMemory inspection_result 字段", "如 FAIL 可触发重检流程"],
        "color": ORANGE,
    },
]

for step in step_details:
    slide = add_slide()
    color = step["color"]

    # Header
    slide_header(slide, f"步骤 {step['num']}  ·  {step['name']}",
                 subtitle=f"技能序列：{step['seq']}", accent_color=color)

    # Description box
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.7), BG_CARD)
    add_textbox(slide, step["desc"],
                Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.65),
                font_size=13, color=LIGHT_GRAY, italic=True)

    # Skill cards
    skills = step["skills"]
    n = len(skills)
    card_w = Inches(12.5 / n) if n <= 6 else Inches(2.05)
    for si, (sid, fname, sdesc, scolor) in enumerate(skills):
        lx = Inches(0.35 + si * (12.6 / n))
        ty = Inches(2.05)
        add_rect(slide, lx, ty, card_w - Inches(0.08), Inches(2.5), BG_CARD)
        add_rect(slide, lx, ty, card_w - Inches(0.08), Inches(0.35), scolor)
        add_textbox(slide, sid,   lx + Inches(0.06), ty + Inches(0.03),
                    Inches(0.6), Inches(0.3), font_size=12, bold=True, color=BG_DARK)
        add_textbox(slide, fname, lx + Inches(0.06), ty + Inches(0.42),
                    card_w - Inches(0.18), Inches(0.4),
                    font_size=11, bold=True, color=scolor)
        add_textbox(slide, sdesc, lx + Inches(0.06), ty + Inches(0.85),
                    card_w - Inches(0.18), Inches(1.5),
                    font_size=10, color=LIGHT_GRAY, wrap=True)
        # Arrow connector
        if si < n - 1:
            ax = lx + card_w - Inches(0.06)
            add_textbox(slide, "→", ax - Inches(0.15), ty + Inches(0.9), Inches(0.3), Inches(0.3),
                        font_size=14, color=color, align=PP_ALIGN.CENTER, bold=True)

    # Notes
    add_rect(slide, Inches(0.3), Inches(4.72), Inches(12.7), Inches(0.35), ACCENT)
    add_textbox(slide, "⚙  设计要点", Inches(0.45), Inches(4.73), Inches(2.0), Inches(0.3),
                font_size=11, bold=True, color=HIGHLIGHT)
    notes_text = "    ·  ".join(step["notes"])
    add_textbox(slide, "·  " + notes_text,
                Inches(0.45), Inches(5.12), Inches(12.2), Inches(0.5),
                font_size=11, color=LIGHT_GRAY)

    # Sequence bar at bottom
    add_rect(slide, Inches(0.3), Inches(5.75), Inches(12.7), Inches(1.55), BG_CARD)
    add_textbox(slide, "完整技能调用序列", Inches(0.45), Inches(5.8), Inches(3.0), Inches(0.3),
                font_size=10, bold=True, color=color)
    add_textbox(slide, step["seq"],
                Inches(0.45), Inches(6.15), Inches(12.2), Inches(0.9),
                font_size=15, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Key Design Constraints
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "关键设计约束", "硬件精度要求与控制参数", accent_color=ORANGE)

constraints = [
    ("⚡", "力控精度", "0.5 – 2 N",
     "GR2 灵巧手亚牛顿级力控，避免损伤镜头及镜头帽", ORANGE),
    ("🎯", "对准精度", "± 0.3 mm",
     "镜头帽插拔、工装槽位放置、周转盒码放均需达到此重复定位精度", HIGHLIGHT),
    ("🤝", "双臂同步", "< 100 ms",
     "M2 bimanual_sync_move 及 O2 press_dual_buttons 时间偏差上限", GREEN),
    ("👁", "视觉伺服", "实时",
     "所有抓取操作融合 P1 visual_localize 实时补偿位置误差", YELLOW),
    ("🛡", "防碰撞", "M4 plan_path",
     "轨迹规划需考虑检测台、周转盒、扫码器等障碍物", RGBColor(0xC7,0x77,0xFF)),
]

for ci, (icon, title, value, desc, color) in enumerate(constraints):
    lx = Inches(0.35 + (ci % 3) * 4.3)
    ty = Inches(1.3 + (ci // 3) * 2.5)
    add_rect(slide, lx, ty, Inches(4.1), Inches(2.2), BG_CARD)
    add_rect(slide, lx, ty, Inches(4.1), Inches(0.1), color)
    add_textbox(slide, icon,  lx + Inches(0.15), ty + Inches(0.18), Inches(0.55), Inches(0.55),
                font_size=28, color=color)
    add_textbox(slide, title, lx + Inches(0.75), ty + Inches(0.18), Inches(3.2), Inches(0.35),
                font_size=14, bold=True, color=color)
    add_textbox(slide, value, lx + Inches(0.75), ty + Inches(0.55), Inches(3.2), Inches(0.45),
                font_size=20, bold=True, color=WHITE)
    add_textbox(slide, desc,  lx + Inches(0.15), ty + Inches(1.1),  Inches(3.8), Inches(0.9),
                font_size=11, color=LIGHT_GRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — RoboOS File Mapping
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "RoboOS 适配文件一览", "本次改动的所有文件及其作用", accent_color=GREEN)

files = [
    ("新建", "slaver/gr2_robot_local/skill.py",  "20 项原子技能（stdio 本地模式）",      GREEN),
    ("新建", "slaver/gr2_robot_remote/skill.py", "20 项原子技能（HTTP 远程模式）",       GREEN),
    ("修改", "master/scene/profile.yaml",         "场景替换为气密检测工作台",             HIGHLIGHT),
    ("修改", "master/agents/prompts.py",          "注入 GR2 平台上下文 & 9步参考流程",   HIGHLIGHT),
    ("修改", "slaver/config.yaml",               "机器人指向 gr2_robot / max_tools=5", ORANGE),
    ("修改", "slaver/tools/memory.py",           "新增 none 动作类型（感知/等待类技能）", ORANGE),
]

for ri, (tag, path, desc, color) in enumerate(files):
    ty = Inches(1.3 + ri * 0.93)
    add_rect(slide, Inches(0.3),   ty, Inches(0.7),  Inches(0.75), color)
    add_textbox(slide, tag, Inches(0.3), ty + Inches(0.18), Inches(0.7), Inches(0.35),
                font_size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.1),   ty, Inches(11.9), Inches(0.75), BG_CARD)
    add_textbox(slide, path, Inches(1.2), ty + Inches(0.06), Inches(6.0), Inches(0.35),
                font_size=13, bold=True, color=color)
    add_textbox(slide, desc, Inches(1.2), ty + Inches(0.4),  Inches(11.5), Inches(0.3),
                font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
OUT = "/home/haoanw/workspace/RoboOS/doc/RoboOS_Hikvision_Workflow.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
